import os
import time
from typing import Any, Dict, List, Optional

from pydantic import BaseModel, Field

from backend.config import settings
from backend.tools.base import BaseTool


class PPTToolInput(BaseModel):
    action: str = Field(
        description="PPT action, such as create_slides, update_slide, get_slides, export_markdown, or export_pdf."
    )
    task_id: Optional[str] = None
    slides: Optional[List[Dict[str, Any]]] = None
    title: Optional[str] = None
    slide_id: Optional[str] = None
    deck_spec: Optional[Dict[str, Any]] = None


class PPTTool(BaseTool):
    def __init__(self, mock_mode: bool = None):
        super().__init__("PPTTool", mock_mode if mock_mode is not None else settings.MOCK_MODE)
        self._output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "data", "slides")
        os.makedirs(self._output_dir, exist_ok=True)

    def _build_langchain_tool(self):
        from langchain_core.tools import StructuredTool

        return StructuredTool.from_function(
            name="ppt_tool",
            description="Create, update, export, or read an Agent-Pilot PPT artifact.",
            coroutine=self._run,
            args_schema=PPTToolInput,
        )

    def _download_url(self, filepath: str) -> str:
        filename = os.path.basename(filepath)
        return f"/api/files/slides/{filename}"

    async def _run(
        self,
        action: str,
        task_id: str = None,
        slides: List[Dict] = None,
        title: str = None,
        slide_id: str = None,
        deck_spec: Dict = None,
    ) -> Dict[str, Any]:
        self._log("info", f"PPT action: {action}", {"task_id": task_id})
        if not self._validate_input({"action": action}, ["action"]):
            return {"success": False, "error": "Missing required parameter: action"}

        if self.mock_mode:
            return await self._mock_run(action, task_id, slides, title, slide_id)
        return await self._real_run(action, task_id, slides, title, slide_id, deck_spec)

    async def _mock_run(
        self,
        action: str,
        task_id: str,
        slides: List[Dict],
        title: str,
        slide_id: str,
    ) -> Dict[str, Any]:
        await self._simulate_delay(0.5)
        self._log("info", f"Mock PPT action executed: {action}")
        mock_responses = {
            "create_slides": {
                "success": True,
                "slide_id": f"slide_{int(time.time() * 1000)}",
                "task_id": task_id,
                "title": title or "Presentation",
                "slides_count": len(slides) if slides else 5,
                "file_path": None,
                "download_url": None,
            },
            "update_slide": {"success": True, "slide_id": slide_id, "updated": True},
            "get_slides": {"success": True, "slides": slides or self._generate_mock_slides()},
            "export_markdown": {"success": True, "file_path": f"{self._output_dir}/mock_{int(time.time() * 1000)}.md"},
            "export_pdf": {"success": True, "file_path": f"{self._output_dir}/mock_{int(time.time() * 1000)}.pdf"},
        }
        return mock_responses.get(action, {"success": True, "action": action})

    async def _real_run(
        self,
        action: str,
        task_id: str,
        slides: List[Dict],
        title: str,
        slide_id: str,
        deck_spec: Dict = None,
    ) -> Dict[str, Any]:
        try:
            if action == "create_slides":
                return await self._create_pptx(task_id, title, slides, deck_spec)
            if action == "update_slide":
                return await self._update_slide(slide_id, slides)
            if action == "get_slides":
                return await self._get_slides(task_id)
            if action == "export_markdown":
                return await self._export_markdown(title, slides, deck_spec)
            if action == "export_pdf":
                return await self._export_pdf(title, slides, deck_spec)
            return {"success": False, "error": f"Unknown action: {action}"}
        except Exception as exc:
            self._log("error", f"PPT action failed: {str(exc)}")
            return {"success": False, "error": str(exc)}

    async def _create_pptx(
        self,
        task_id: str,
        title: str,
        slides: List[Dict],
        deck_spec: Dict = None,
    ) -> Dict[str, Any]:
        self._log("info", f"Creating plain PPT for task {task_id}")
        try:
            filename = f"{task_id}_{int(time.time() * 1000)}.pptx"
            filepath = os.path.join(self._output_dir, filename)
            payload_slides = self._normalize_slides_payload(slides, deck_spec)
            deck_title = (deck_spec or {}).get("title") or title or "Presentation"
            self._write_basic_pptx(filepath, deck_title, payload_slides)
            return {
                "success": True,
                "slide_id": f"slide_{int(time.time() * 1000)}",
                "task_id": task_id,
                "title": deck_title,
                "slides_count": len(payload_slides),
                "deck_spec": {"title": deck_title, "slides": payload_slides},
                "file_path": filepath,
                "download_url": self._download_url(filepath),
            }
        except ImportError:
            self._log("warning", "python-pptx not available")
            return {
                "success": False,
                "error": "python-pptx not installed. Run: python -m pip install python-pptx",
            }
        except Exception as exc:
            self._log("error", f"PPTX creation failed: {str(exc)}")
            return {"success": False, "error": str(exc)}

    async def _render_with_deck_spec(self, deck_spec: Dict, task_id: str) -> Dict[str, Any]:
        return await self._create_pptx(task_id, (deck_spec or {}).get("title"), [], deck_spec)

    @staticmethod
    def _normalize_slides_payload(slides: Optional[List[Dict]], deck_spec: Optional[Dict]) -> List[Dict[str, Any]]:
        source = []
        if isinstance(deck_spec, dict):
            source = deck_spec.get("slides") or []
        if not source:
            source = slides or []

        normalized: List[Dict[str, Any]] = []
        for index, item in enumerate(source):
            if not isinstance(item, dict):
                continue
            bullets = item.get("bullets")
            if not isinstance(bullets, list):
                bullets = []
            if not bullets:
                content = item.get("content")
                if isinstance(content, str):
                    bullets = [line.strip() for line in content.splitlines() if line.strip()]
                elif isinstance(content, list):
                    bullets = [str(v).strip() for v in content if str(v).strip()]
                elif isinstance(content, dict):
                    bullets = [str(v).strip() for v in content.values() if str(v).strip()]
            normalized.append(
                {
                    "index": index,
                    "title": str(item.get("title") or f"Slide {index + 1}"),
                    "bullets": bullets[:10],
                }
            )

        if not normalized:
            return [{"index": 0, "title": "Presentation", "bullets": []}]
        return normalized

    @staticmethod
    def _write_basic_pptx(filepath: str, title: str, slides: List[Dict[str, Any]]) -> None:
        from pptx import Presentation

        prs = Presentation()
        title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[6]
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

        first = slides[0] if slides else {"title": title, "bullets": []}
        cover = prs.slides.add_slide(title_layout)
        if cover.shapes.title:
            cover.shapes.title.text = str(first.get("title") or title)
        if len(cover.placeholders) > 1 and hasattr(cover.placeholders[1], "text"):
            cover.placeholders[1].text = ""

        for slide in slides[1:] if len(slides) > 1 else []:
            page = prs.slides.add_slide(content_layout)
            if page.shapes.title:
                page.shapes.title.text = str(slide.get("title") or "")
            body = page.shapes.placeholders[1] if len(page.shapes.placeholders) > 1 else None
            if body and hasattr(body, "text_frame"):
                tf = body.text_frame
                tf.clear()
                bullets = slide.get("bullets") or []
                if bullets:
                    tf.paragraphs[0].text = str(bullets[0])
                    for item in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = str(item)
                        p.level = 0

        prs.save(filepath)

    async def _export_markdown(self, title: str, slides: List[Dict], deck_spec: Dict = None) -> Dict[str, Any]:
        try:
            from backend.services.deck_renderer import SlidevRenderer
            from backend.services.deck_spec import DeckSpec

            if deck_spec:
                deck = DeckSpec.from_dict(deck_spec)
            else:
                deck = DeckSpec(title=title or "Presentation")
                for slide in (slides or self._generate_mock_slides()):
                    deck.add_slide(
                        title=slide.get("title", ""),
                        layout=slide.get("layout", "content"),
                        bullets=slide.get("bullets", []),
                    )
            renderer = SlidevRenderer(self._output_dir)
            filepath = renderer.render_to_file(deck)
            return {"success": True, "file_path": filepath}
        except Exception as exc:
            self._log("error", f"Markdown export failed: {str(exc)}")
            return {"success": False, "error": str(exc)}

    async def _export_pdf(self, title: str, slides: List[Dict], deck_spec: Dict = None) -> Dict[str, Any]:
        try:
            from backend.services.deck_renderer import MarkdownToPdfRenderer, SlidevRenderer
            from backend.services.deck_spec import DeckSpec

            if deck_spec:
                deck = DeckSpec.from_dict(deck_spec)
            else:
                deck = DeckSpec(title=title or "Presentation")
                for slide in (slides or self._generate_mock_slides()):
                    deck.add_slide(title=slide.get("title", ""), layout=slide.get("layout", "content"))

            md_renderer = SlidevRenderer(self._output_dir)
            md_path = md_renderer.render_to_file(deck)
            pdf_renderer = MarkdownToPdfRenderer()
            return await pdf_renderer.render(md_path)
        except Exception as exc:
            self._log("error", f"PDF export failed: {str(exc)}")
            return {"success": False, "error": str(exc)}

    async def _update_slide(self, slide_id: str, slides: List[Dict]) -> Dict[str, Any]:
        self._log("info", f"Updating slide {slide_id}")
        return {"success": True, "slide_id": slide_id, "updated": True}

    async def _get_slides(self, task_id: str) -> Dict[str, Any]:
        self._log("info", f"Getting slides for task {task_id}")
        return {"success": True, "slides": []}

    def _generate_mock_slides(self) -> List[Dict]:
        return [
            {"index": 0, "title": "封面", "layout": "title", "bullets": []},
            {"index": 1, "title": "背景与目标", "layout": "content", "bullets": ["目标 1", "目标 2"]},
            {"index": 2, "title": "核心方案", "layout": "content", "bullets": ["方案 A", "方案 B"]},
            {"index": 3, "title": "风险与待办", "layout": "content", "bullets": ["风险 1", "待办 1"]},
            {"index": 4, "title": "下一步", "layout": "content", "bullets": ["行动项"]},
        ]

    async def _simulate_delay(self, seconds: float):
        import asyncio

        await asyncio.sleep(seconds)
