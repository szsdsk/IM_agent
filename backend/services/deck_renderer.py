from __future__ import annotations

"""DeckSpec renderers.

The PPTX renderer intentionally keeps layout decisions deterministic: the LLM
provides content and intent, while this module draws stable visual structures.
"""

import logging
import os
from typing import Any, Dict

from backend.services.deck_spec import DeckSpec, SlideSpec, ThemeConfig

logger = logging.getLogger(__name__)


class SlidevRenderer:
    """Render DeckSpec into simple Slidev markdown."""

    def __init__(self, output_dir: str = None):
        self.output_dir = output_dir or "./data/slides"
        os.makedirs(self.output_dir, exist_ok=True)

    def render(self, deck: DeckSpec) -> str:
        lines = [
            "---",
            f"title: {deck.title}",
            "author: Agent-Pilot",
            "transition: fade",
            "theme: default",
            "---",
            "",
        ]
        for slide in deck.slides:
            lines.extend(self._render_slide(slide))
            lines.append("")
        return "\n".join(lines)

    def _render_slide(self, slide: SlideSpec) -> list[str]:
        lines = [f"## {slide.title}", ""]
        for bullet in slide.bullets:
            lines.append(f"- {bullet}")
        if slide.speaker_notes:
            lines.extend(["", f"// {slide.speaker_notes}"])
        return lines

    def render_to_file(self, deck: DeckSpec, filename: str = None) -> str:
        if filename is None:
            filename = f"{deck.title.replace(' ', '_')}_{int(__import__('time').time())}.md"
        path = os.path.join(self.output_dir, filename)
        with open(path, "w", encoding="utf-8") as f:
            f.write(self.render(deck))
        logger.info("Slidev markdown saved to %s", path)
        return path


class PptxGenRenderer:
    """Render normalized DeckSpec into a richer PPTX file with python-pptx."""

    def __init__(self, output_dir: str = None):
        self.output_dir = output_dir or "./data/slides"
        os.makedirs(self.output_dir, exist_ok=True)

    @staticmethod
    def _hex_to_rgb(hex_color: str):
        from pptx.dml.color import RGBColor

        h = (hex_color or "#111827").lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _safe_text(value: Any, limit: int = 120) -> str:
        if value is None:
            return ""
        text = str(value).replace("\n", " ").strip()
        return text if len(text) <= limit else f"{text[:limit - 1]}..."

    async def render(self, deck: DeckSpec, filename: str = None) -> Dict[str, Any]:
        if filename is None:
            filename = f"{deck.title.replace(' ', '_')}.pptx"
        filepath = os.path.join(self.output_dir, filename)

        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.warning("python-pptx not installed")
            return {"success": False, "error": "python-pptx not installed"}

        theme = ThemeConfig.get_theme(deck.theme)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_spec in deck.slides:
            self._add_slide(prs, slide_spec, theme)

        prs.save(filepath)
        return {
            "success": True,
            "filepath": filepath,
            "slides_count": len(deck.slides),
            "title": deck.title,
        }

    def _add_slide(self, prs, slide: SlideSpec, theme: ThemeConfig) -> None:
        layout = prs.slide_layouts[6]
        slide_obj = prs.slides.add_slide(layout)
        self._apply_slide_background(slide_obj, theme, slide.index + 1)

        if slide.layout in {"title", "hero"}:
            self._render_hero_layout(slide_obj, slide, theme)
        elif slide.layout == "section_divider":
            self._render_section_divider_layout(slide_obj, slide, theme)
        elif slide.layout == "metrics":
            self._render_metrics_layout(slide_obj, slide, theme)
        elif slide.layout == "timeline":
            self._render_timeline_layout(slide_obj, slide, theme)
        elif slide.layout in {"two_column", "comparison"}:
            self._render_comparison_layout(slide_obj, slide, theme)
        elif slide.layout in {"diagram", "process"}:
            self._render_process_layout(slide_obj, slide, theme)
        elif slide.layout == "cards":
            self._render_cards_layout(slide_obj, slide, theme)
        elif slide.layout == "chart":
            self._render_chart_layout(slide_obj, slide, theme)
        elif slide.layout == "closing":
            self._render_closing_layout(slide_obj, slide, theme)
        else:
            self._render_content_layout(slide_obj, slide, theme)

        if slide.speaker_notes:
            slide_obj.notes_slide.notes_text_frame.text = slide.speaker_notes

    def _apply_slide_background(self, slide_obj, theme: ThemeConfig, page_number: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        colors = theme.colors
        theme_name = theme.name
        fill = slide_obj.background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(colors.bg_white)

        # 每套主题使用不同的底纹语言，避免所有页面都像同一张卡片模板。
        if theme_name == "entertainment":
            for x, y, size, color in [
                (10.35, -1.3, 3.8, colors.primary_light),
                (9.25, 4.85, 2.2, colors.accent),
                (-0.95, 5.25, 2.1, colors.primary_light),
            ]:
                orb = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
                orb.fill.solid()
                orb.fill.fore_color.rgb = self._hex_to_rgb(color)
                orb.line.fill.background()
            ribbon = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.18), Inches(13.333), Inches(0.14))
            ribbon.rotation = -4
            ribbon.fill.solid()
            ribbon.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
            ribbon.line.fill.background()
        elif theme_name == "tech_dark":
            rail = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
            rail.fill.solid()
            rail.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
            rail.line.fill.background()
            for i in range(4):
                stripe = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.8 + i * 0.55), Inches(0.5 + i * 0.45), Inches(2.2), Inches(0.04))
                stripe.fill.solid()
                stripe.fill.fore_color.rgb = self._hex_to_rgb(colors.divider)
                stripe.line.fill.background()
        else:
            accent = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.75), Inches(-1.25), Inches(3.8), Inches(3.8))
            accent.fill.solid()
            accent.fill.fore_color.rgb = self._hex_to_rgb(colors.primary_light)
            accent.line.fill.background()

        bar = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(colors.primary)
        bar.line.fill.background()

        page_box = slide_obj.shapes.add_textbox(Inches(12.45), Inches(7.12), Inches(0.75), Inches(0.3))
        p = page_box.text_frame.paragraphs[0]
        p.text = str(page_number)
        p.font.size = Pt(10)
        p.font.color.rgb = self._hex_to_rgb(colors.text_light)
        p.alignment = PP_ALIGN.RIGHT

    def _add_text(
        self,
        slide_obj,
        text: str,
        x,
        y,
        w,
        h,
        theme: ThemeConfig,
        size: int = 16,
        bold: bool = False,
        color: str = None,
        align=None,
    ):
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Pt

        box = slide_obj.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = theme.font_body
        p.font.color.rgb = self._hex_to_rgb(color or theme.colors.text_dark)
        p.alignment = align or PP_ALIGN.LEFT
        return box

    def _add_card(self, slide_obj, x, y, w, h, theme: ThemeConfig, fill: str = None, line: str = None):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide_obj.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(fill or theme.colors.bg_light)
        shape.line.color.rgb = self._hex_to_rgb(line or theme.colors.divider)
        return shape

    def _add_rule(self, slide_obj, x, y, w, h, color: str):
        from pptx.enum.shapes import MSO_SHAPE

        rule = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        rule.fill.solid()
        rule.fill.fore_color.rgb = self._hex_to_rgb(color)
        rule.line.fill.background()
        return rule

    def _add_title_header(self, slide_obj, slide: SlideSpec, theme: ThemeConfig, eyebrow: str) -> None:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        self._add_text(slide_obj, eyebrow, Inches(0.78), Inches(0.32), Inches(3.2), Inches(0.26), theme, 9, True, colors.accent)
        self._add_text(slide_obj, self._safe_text(slide.title, 48), Inches(0.75), Inches(0.62), Inches(10.8), Inches(0.62), theme, 28, True, colors.text_dark, PP_ALIGN.LEFT)
        self._add_rule(slide_obj, Inches(0.75), Inches(1.38), Inches(1.35), Inches(0.06), colors.accent)

    def _slide_bullets(self, slide: SlideSpec, limit: int = 6) -> list[str]:
        bullets = [self._safe_text(item, 88) for item in (slide.bullets or []) if str(item).strip()]
        if not bullets and slide.content:
            if isinstance(slide.content, str):
                bullets = [self._safe_text(line.strip("-• "), 88) for line in slide.content.splitlines() if line.strip()]
            elif isinstance(slide.content, dict):
                bullets = [self._safe_text(value, 88) for value in slide.content.values() if value]
            elif isinstance(slide.content, list):
                bullets = [self._safe_text(value, 88) for value in slide.content if value]
        return bullets[:limit]

    def _render_hero_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        if theme.name == "entertainment":
            stage = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.1))
            stage.fill.solid()
            stage.fill.fore_color.rgb = self._hex_to_rgb(colors.bg_dark)
            stage.line.fill.background()

            moon = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.75), Inches(0.55), Inches(3.3), Inches(3.3))
            moon.fill.solid()
            moon.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
            moon.line.fill.background()
            halo = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.45), Inches(3.95), Inches(4.6), Inches(4.6))
            halo.fill.solid()
            halo.fill.fore_color.rgb = self._hex_to_rgb(colors.primary)
            halo.line.fill.background()

            self._add_rule(slide_obj, Inches(0.82), Inches(0.82), Inches(0.12), Inches(5.35), colors.accent)
            self._add_text(slide_obj, "ENTERTAINMENT DECK", Inches(1.18), Inches(0.95), Inches(4.6), Inches(0.34), theme, 10, True, colors.accent)
            self._add_text(slide_obj, self._safe_text(slide.title, 52), Inches(1.08), Inches(1.65), Inches(7.45), Inches(2.0), theme, 42, True, colors.text_light)
            bullets = self._slide_bullets(slide, 3)
            subtitle = " / ".join(bullets) if bullets else self._safe_text(slide.content, 88)
            if subtitle:
                self._add_text(slide_obj, subtitle, Inches(1.12), Inches(4.22), Inches(6.9), Inches(0.75), theme, 16, False, colors.primary_light)
            for i, item in enumerate(bullets[:3]):
                y = 5.35 + i * 0.36
                self._add_rule(slide_obj, Inches(8.92), Inches(y), Inches(0.48), Inches(0.06), colors.accent)
                self._add_text(slide_obj, self._safe_text(item, 24), Inches(9.55), Inches(y - 0.1), Inches(2.4), Inches(0.24), theme, 10, True, colors.text_light)
            return

        hero_bg = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.1))
        hero_bg.fill.solid()
        hero_bg.fill.fore_color.rgb = self._hex_to_rgb(colors.primary)
        hero_bg.line.fill.background()

        glow = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.7), Inches(0.55), Inches(4.2), Inches(4.2))
        glow.fill.solid()
        glow.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
        glow.line.fill.background()

        self._add_rule(slide_obj, Inches(0.82), Inches(0.85), Inches(0.1), Inches(5.1), colors.accent)

        self._add_text(slide_obj, (slide.visual_profile or "Agent-Pilot").upper(), Inches(1.28), Inches(1.24), Inches(4.0), Inches(0.3), theme, 10, True, colors.accent)
        self._add_text(slide_obj, self._safe_text(slide.title, 56), Inches(1.23), Inches(1.72), Inches(7.8), Inches(1.8), theme, 40, True, colors.text_light)
        bullets = self._slide_bullets(slide, 3)
        subtitle = " | ".join(bullets) if bullets else self._safe_text(slide.content, 88)
        if subtitle:
            self._add_text(slide_obj, subtitle, Inches(1.3), Inches(3.75), Inches(7.4), Inches(0.72), theme, 17, False, colors.primary_light)
        for i, item in enumerate(bullets[:3]):
            x = Inches(1.25 + i * 2.35)
            self._add_card(slide_obj, x, Inches(5.0), Inches(2.05), Inches(0.65), theme, colors.primary_light, colors.primary_light)
            self._add_text(slide_obj, self._safe_text(item, 18), x + Inches(0.12), Inches(5.18), Inches(1.8), Inches(0.24), theme, 10, True, colors.primary, PP_ALIGN.CENTER)

    def _render_content_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        self._add_title_header(slide_obj, slide, theme, "KEY POINTS")
        self._add_rule(slide_obj, Inches(0.88), Inches(1.78), Inches(0.08), Inches(4.75), colors.primary)
        for i, bullet in enumerate(self._slide_bullets(slide, 6)):
            y = Inches(1.85 + i * 0.78)
            dot = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.72), y + Inches(0.12), Inches(0.38), Inches(0.38))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._hex_to_rgb(colors.accent if i == 0 else colors.primary_light)
            dot.line.fill.background()
            self._add_text(slide_obj, f"{i + 1:02d}", Inches(1.28), y + Inches(0.06), Inches(0.52), Inches(0.25), theme, 10, True, colors.accent, PP_ALIGN.LEFT)
            self._add_text(slide_obj, bullet, Inches(1.92), y, Inches(9.75), Inches(0.48), theme, 16 if i == 0 else 14, i == 0, colors.text_dark)

    def _render_section_divider_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        colors = theme.colors
        banner = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(-0.35), Inches(1.45), Inches(14.2), Inches(2.4))
        banner.rotation = -3
        banner.fill.solid()
        banner.fill.fore_color.rgb = self._hex_to_rgb(colors.primary)
        banner.line.fill.background()
        side = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.55), Inches(0.92), Inches(2.6), Inches(2.6))
        side.fill.solid()
        side.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
        side.line.fill.background()
        self._add_text(slide_obj, self._safe_text(slide.title, 48), Inches(1.15), Inches(2.08), Inches(8.8), Inches(0.9), theme, 36, True, colors.text_light)
        for i, bullet in enumerate(self._slide_bullets(slide, 4)):
            self._add_rule(slide_obj, Inches(1.22), Inches(4.35 + i * 0.48), Inches(0.42), Inches(0.05), colors.accent)
            self._add_text(slide_obj, bullet, Inches(1.85), Inches(4.22 + i * 0.48), Inches(9.4), Inches(0.32), theme, 14, False, colors.text_dark)

    def _render_metrics_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        self._add_title_header(slide_obj, slide, theme, "METRICS")
        metrics = slide.highlight_metrics or []
        if not metrics:
            self._render_cards_layout(slide_obj, slide, theme)
            return
        self._add_rule(slide_obj, Inches(0.95), Inches(3.48), Inches(11.25), Inches(0.05), colors.divider)
        positions = [(0.95, 2.05), (4.45, 2.05), (7.95, 2.05), (2.7, 4.35), (6.45, 4.35)]
        for index, metric in enumerate(metrics[:5]):
            x, y = positions[index]
            if index < 3:
                self._add_rule(slide_obj, Inches(x), Inches(y + 1.02), Inches(2.35), Inches(0.06), colors.accent)
            self._add_text(slide_obj, self._safe_text(metric.get("value"), 14), Inches(x), Inches(y), Inches(2.75), Inches(0.62), theme, 34, True, colors.primary, PP_ALIGN.LEFT)
            self._add_text(slide_obj, self._safe_text(metric.get("label"), 42), Inches(x + 0.03), Inches(y + 0.82), Inches(2.65), Inches(0.5), theme, 12, False, colors.text_muted, PP_ALIGN.LEFT)

    def _render_timeline_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        self._add_title_header(slide_obj, slide, theme, "TIMELINE")
        items = slide.timeline or [{"label": f"Phase {i + 1}", "text": item} for i, item in enumerate(self._slide_bullets(slide, 5))]
        if not items:
            self._render_content_layout(slide_obj, slide, theme)
            return
        start_x, y = 1.05, 3.05
        gap = 10.9 / max(len(items[:5]) - 1, 1)
        line = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(start_x), Inches(y + 0.25), Inches(10.9), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(colors.primary_light)
        line.line.fill.background()
        for index, item in enumerate(items[:5]):
            x = start_x + index * gap
            dot = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y - 0.1), Inches(0.65), Inches(0.65))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._hex_to_rgb(colors.primary if index % 2 == 0 else colors.accent)
            dot.line.fill.background()
            self._add_text(slide_obj, f"{index + 1}", Inches(x + 0.15), Inches(y + 0.07), Inches(0.32), Inches(0.2), theme, 11, True, colors.text_light, PP_ALIGN.CENTER)
            label_y = y - 0.82 if index % 2 == 0 else y + 0.85
            body_y = y - 0.42 if index % 2 == 0 else y + 1.25
            self._add_text(slide_obj, self._safe_text(item.get("label"), 14), Inches(x - 0.35), Inches(label_y), Inches(1.45), Inches(0.32), theme, 12, True, colors.primary, PP_ALIGN.CENTER)
            self._add_text(slide_obj, self._safe_text(item.get("text"), 42), Inches(x - 0.65), Inches(body_y), Inches(2.0), Inches(0.5), theme, 10, False, colors.text_dark, PP_ALIGN.CENTER)

    def _render_comparison_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.util import Inches

        colors = theme.colors
        self._add_title_header(slide_obj, slide, theme, "COMPARISON")
        bullets = self._slide_bullets(slide, 6)
        mid = max(1, len(bullets) // 2)
        groups = [
            ("Current / Pain Points", bullets[:mid], colors.primary_light),
            ("Target / Solution", bullets[mid:], colors.bg_light),
        ]
        for index, (heading, items, fill) in enumerate(groups):
            x = 0.85 + index * 6.1
            block = self._add_card(slide_obj, Inches(x), Inches(1.85), Inches(5.55), Inches(4.65), theme, fill)
            block.shadow.inherit = False
            self._add_rule(slide_obj, Inches(x + 0.35), Inches(2.55), Inches(1.1), Inches(0.06), colors.accent)
            self._add_text(slide_obj, heading, Inches(x + 0.35), Inches(2.12), Inches(4.6), Inches(0.32), theme, 17, True, colors.primary)
            for i, item in enumerate(items[:4]):
                self._add_text(slide_obj, f"- {item}", Inches(x + 0.38), Inches(2.75 + i * 0.72), Inches(4.75), Inches(0.35), theme, 13, False, colors.text_dark)

    def _render_process_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        self._add_title_header(slide_obj, slide, theme, "PROCESS")
        steps = slide.process_steps or [{"label": f"{i + 1:02d}", "text": item} for i, item in enumerate(self._slide_bullets(slide, 6))]
        if not steps:
            steps = [{"label": "01", "text": "Diagram area"}]
        self._add_rule(slide_obj, Inches(1.25), Inches(2.55), Inches(10.4), Inches(0.05), colors.divider)
        for index, step in enumerate(steps[:6]):
            row = index // 3
            col = index % 3
            x = 1.0 + col * 4.05
            y = 2.0 + row * 2.0
            badge = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.25), Inches(0.55), Inches(0.55))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self._hex_to_rgb(colors.primary)
            badge.line.fill.background()
            self._add_text(slide_obj, self._safe_text(step.get("label"), 4), Inches(x + 0.2), Inches(y + 0.38), Inches(0.5), Inches(0.2), theme, 9, True, colors.text_light, PP_ALIGN.CENTER)
            self._add_text(slide_obj, self._safe_text(step.get("text"), 42), Inches(x + 0.9), Inches(y + 0.28), Inches(2.3), Inches(0.52), theme, 13, index == 0, colors.text_dark)
            if index < min(len(steps), 6) - 1 and col < 2:
                connector = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 3.3), Inches(y + 0.55), Inches(0.45), Inches(0.05))
                connector.fill.solid()
                connector.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
                connector.line.fill.background()

    def _render_cards_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.util import Inches

        colors = theme.colors
        self._add_title_header(slide_obj, slide, theme, "INSIGHTS")
        sections = slide.sections or [{"title": item[:12], "body": item} for item in self._slide_bullets(slide, 6)]
        for index, item in enumerate(sections[:6]):
            row = index // 3
            col = index % 3
            x = 0.95 + col * 4.05
            y = 1.82 + row * 2.05 + (0.25 if col == 1 else 0)
            self._add_text(slide_obj, f"{index + 1:02d}", Inches(x), Inches(y), Inches(0.7), Inches(0.28), theme, 12, True, colors.accent)
            self._add_rule(slide_obj, Inches(x), Inches(y + 0.42), Inches(2.95), Inches(0.04), colors.primary_light)
            self._add_text(slide_obj, self._safe_text(item.get("title"), 18), Inches(x), Inches(y + 0.62), Inches(2.95), Inches(0.3), theme, 15, True, colors.primary)
            self._add_text(slide_obj, self._safe_text(item.get("body"), 64), Inches(x), Inches(y + 1.04), Inches(3.1), Inches(0.52), theme, 11, False, colors.text_dark)

    def _render_closing_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        colors = theme.colors
        panel = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.1))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._hex_to_rgb(colors.primary)
        panel.line.fill.background()
        slash = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.9), Inches(-0.2), Inches(0.75), Inches(8.0))
        slash.rotation = 14
        slash.fill.solid()
        slash.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
        slash.line.fill.background()
        self._add_text(slide_obj, self._safe_text(slide.title, 34), Inches(2.0), Inches(2.0), Inches(9.3), Inches(0.8), theme, 34, True, colors.text_light, PP_ALIGN.CENTER)
        bullets = self._slide_bullets(slide, 3)
        text = " | ".join(bullets) if bullets else "Thanks. Q&A"
        self._add_text(slide_obj, text, Inches(2.15), Inches(3.25), Inches(9.0), Inches(0.55), theme, 16, False, colors.text_light, PP_ALIGN.CENTER)
        mark = slide_obj.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.95), Inches(4.35), Inches(1.35), Inches(1.35))
        mark.fill.solid()
        mark.fill.fore_color.rgb = self._hex_to_rgb(colors.accent)
        mark.line.fill.background()
        self._add_text(slide_obj, "Q&A", Inches(6.05), Inches(4.78), Inches(1.15), Inches(0.3), theme, 18, True, colors.text_light, PP_ALIGN.CENTER)

    def _render_chart_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
        from pptx.chart.data import CategoryChartData

        colors = theme.colors
        chart_data = slide.chart

        if not chart_data or not chart_data.get("series"):
            self._render_content_layout(slide_obj, slide, theme)
            return

        self._add_title_header(slide_obj, slide, theme, "CHART")

        chart_type_str = chart_data.get("type", "bar")
        xl_chart_type = self._chart_type_to_xl(chart_type_str)
        categories = chart_data.get("categories", [])
        series_list = chart_data.get("series", [])

        category_data = CategoryChartData()
        category_data.categories = categories
        for s in series_list:
            values = [float(v) for v in s.get("values", [])]
            category_data.add_series(s.get("name", ""), values)

        chart_frame = slide_obj.shapes.add_chart(
            xl_chart_type,
            Inches(0.75), Inches(1.65),
            Inches(11.8), Inches(5.0),
            category_data,
        )

        chart = chart_frame.chart
        chart.has_legend = len(series_list) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = theme.font_body

        theme_palette = [
            colors.primary, colors.accent, colors.secondary,
            colors.primary_light, colors.divider,
        ]
        plot = chart.plots[0]

        if chart_type_str == "pie":
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.size = Pt(9)
            data_labels.font.name = theme.font_body
            data_labels.font.color.rgb = self._hex_to_rgb(colors.text_dark)
            data_labels.number_format = '0%'
            data_labels.show_percentage = True
            data_labels.show_category_name = True
            data_labels.show_value = False
            # Color each pie slice
            series_format = plot.series[0]
            for i in range(len(categories)):
                point = series_format.points[i]
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = self._hex_to_rgb(theme_palette[i % len(theme_palette)])
        else:
            for i, series in enumerate(plot.series):
                hex_color = theme_palette[i % len(theme_palette)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = self._hex_to_rgb(hex_color)

        if hasattr(plot, 'value_axis') and chart_type_str != "pie":
            value_axis = plot.value_axis
            value_axis.has_title = False
            value_axis.major_gridlines.format.line.color.rgb = self._hex_to_rgb(colors.divider)
            value_axis.format.line.color.rgb = self._hex_to_rgb(colors.divider)
            value_axis.tick_labels.font.size = Pt(9)
            value_axis.tick_labels.font.name = theme.font_body
            value_axis.tick_labels.font.color.rgb = self._hex_to_rgb(colors.text_muted)

        if hasattr(plot, 'category_axis') and chart_type_str != "pie":
            cat_axis = plot.category_axis
            cat_axis.has_title = False
            cat_axis.format.line.color.rgb = self._hex_to_rgb(colors.divider)
            cat_axis.tick_labels.font.size = Pt(9)
            cat_axis.tick_labels.font.name = theme.font_body
            cat_axis.tick_labels.font.color.rgb = self._hex_to_rgb(colors.text_muted)

    @staticmethod
    def _chart_type_to_xl(chart_type: str):
        from pptx.enum.chart import XL_CHART_TYPE
        mapping = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        return mapping.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    def _render_title_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        self._render_hero_layout(slide_obj, slide, theme)

    def _render_two_column_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        self._render_comparison_layout(slide_obj, slide, theme)

    def _render_diagram_layout(self, slide_obj, slide: SlideSpec, theme: ThemeConfig) -> None:
        self._render_process_layout(slide_obj, slide, theme)


class MarkdownToPdfRenderer:
    """Convert Slidev markdown to PDF when Slidev is available."""

    def __init__(self, slidev_path: str = None):
        self.slidev_path = slidev_path or "slidev"

    async def render(self, markdown_path: str, output_path: str = None) -> Dict[str, Any]:
        if output_path is None:
            output_path = markdown_path.replace(".md", ".pdf")
        try:
            import subprocess

            result = subprocess.run(
                [self.slidev_path, "export", markdown_path, "--format", "pdf", "--output", output_path],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0:
                return {"success": True, "filepath": output_path}
            return {"success": False, "error": result.stderr}
        except FileNotFoundError:
            return {"success": False, "error": "Slidev not found"}
        except Exception as e:
            return {"success": False, "error": str(e)}


slidev_renderer = SlidevRenderer()
pptx_renderer = PptxGenRenderer()
pdf_renderer = MarkdownToPdfRenderer()


async def render_deck(deck: DeckSpec, formats: list = None) -> Dict[str, Any]:
    formats = formats or ["pptx"]
    results = {}
    if "markdown" in formats:
        md_path = slidev_renderer.render_to_file(deck)
        results["markdown"] = md_path
    if "pptx" in formats:
        pptx_result = await pptx_renderer.render(deck)
        if pptx_result.get("success"):
            results["pptx"] = pptx_result["filepath"]
    if "pdf" in formats:
        md_path = slidev_renderer.render_to_file(deck)
        pdf_result = await pdf_renderer.render(md_path)
        if pdf_result.get("success"):
            results["pdf"] = pdf_result["filepath"]
    return results
